from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide layout
slide_1 = prs.slides.add_slide(slide_layout)
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]

title.text = "Data Science and Its Uses"
subtitle.text = "An Overview"

# Slide 2: Introduction to Data Science
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide_2 = prs.slides.add_slide(slide_layout)
title = slide_2.shapes.title
content = slide_2.placeholders[1]

title.text = "Introduction to Data Science"
content.text = "Data science is an interdisciplinary field that uses scientific methods, algorithms, processes, and systems to extract knowledge and insights from structured and unstructured data."

# Slide 3: Data Science Process
slide_3 = prs.slides.add_slide(slide_layout)
title = slide_3.shapes.title
content = slide_3.placeholders[1]

title.text = "Data Science Process"
content.text = "1. Data Collection\n2. Data Preprocessing\n3. Data Exploration\n4. Data Modeling\n5. Data Visualization\n6. Communication of Results"

# Slide 4: Applications of Data Science
slide_4 = prs.slides.add_slide(slide_layout)
title = slide_4.shapes.title
content = slide_4.placeholders[1]

title.text = "Applications of Data Science"
content.text = "1. Predictive Analytics\n2. Machine Learning\n3. Natural Language Processing\n4. Image Recognition\n5. Fraud Detection\n6. Recommendation Systems\n7. Healthcare Analytics"

# Slide 5: Conclusion
slide_5 = prs.slides.add_slide(slide_layout)
title = slide_5.shapes.title
content = slide_5.placeholders[1]

title.text = "Conclusion"
content.text = "Data science plays a crucial role in various industries, helping organizations make data-driven decisions, improve processes, and gain competitive advantages."

# Save the presentation
prs.save("Data_Science_Presentation.pptx")

print("Presentation created successfully!")
